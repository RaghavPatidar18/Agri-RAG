import os
import io
import numpy as np
import easyocr
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
from rich.console import Console
from rich.panel import Panel
from rich.text import Text

# Initialize Rich Console for beautiful printing
console = Console()

# Initialize EasyOCR reader (Downloads the model on first run)
with console.status("[bold green]Loading EasyOCR Model (English)..."):
    reader = easyocr.Reader(['en'])

def process_txt(file_path):
    """Reads a plain text file."""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def process_pdf(file_path):
    """Converts PDF to images and applies EasyOCR."""
    extracted_text = []
    
    with console.status(f"[bold cyan]Converting {os.path.basename(file_path)} to images..."):
        # Note: If on Windows, you may need to add poppler_path=r'C:\path\to\poppler\bin' 
        # inside convert_from_path() if it's not in your system PATH.
        pages = convert_from_path(file_path)
        
    for i, page in enumerate(pages):
        with console.status(f"[bold yellow]Running OCR on PDF page {i+1}/{len(pages)}..."):
            # Convert PIL Image to Numpy array for EasyOCR
            img_array = np.array(page)
            # detail=0 returns just the text list, not bounding boxes
            result = reader.readtext(img_array, detail=0)
            
            extracted_text.append(f"--- Page {i+1} ---")
            extracted_text.append(" ".join(result))
            
    return "\n".join(extracted_text)

def process_pptx(file_path):
    """Extracts text from slides and applies OCR to embedded images."""
    prs = Presentation(file_path)
    extracted_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = [f"--- Slide {i+1} ---"]
        
        for shape in slide.shapes:
            # 1. Extract native text
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
                
            # 2. Extract images and run OCR
            if shape.shape_type == 13: # 13 represents a picture
                with console.status(f"[bold magenta]Running OCR on image found in Slide {i+1}..."):
                    image_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    
                    # Convert to RGB (in case of PNG with transparency) and then to Numpy
                    img_array = np.array(img.convert('RGB'))
                    ocr_result = reader.readtext(img_array, detail=0)
                    
                    if ocr_result:
                        slide_text.append("[Image OCR Text]: " + " ".join(ocr_result))
                        
        extracted_text.append("\n".join(slide_text))
        
    return "\n\n".join(extracted_text)

def read_document(file_path):
    """Routes the file to the correct processor and prints it beautifully."""
    if not os.path.exists(file_path):
        console.print(f"[bold red]Error:[/] File '{file_path}' not found.")
        return

    ext = os.path.splitext(file_path)[1].lower()
    content = ""

    try:
        if ext == '.txt':
            content = process_txt(file_path)
        elif ext == '.pdf':
            content = process_pdf(file_path)
        elif ext == '.pptx':
            content = process_pptx(file_path)
        elif ext == '.ppt':
            console.print("[bold red]Warning:[/] Please save your .ppt file as a newer .pptx file first.")
            return
        else:
            console.print(f"[bold red]Unsupported file format:[/] {ext}")
            return
            
        # Print the extracted content beautifully
        panel = Panel(
            Text(content, justify="left", style="white"),
            title=f"[bold cyan]Document Content: {os.path.basename(file_path)}[/]",
            border_style="green",
            padding=(1, 2)
        )
        console.print(panel)

    except Exception as e:
        console.print(f"[bold red]An error occurred while processing {file_path}:[/] {e}")

if __name__ == "__main__":
    sample_file = "/Users/raghavpatidar/Portfolio/blog-agent/documents/CRISIL_Agri_Dashboard_September_2025-28.pdf"
    read_document(sample_file)